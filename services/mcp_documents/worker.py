from __future__ import annotations

import asyncio
import base64
import io
import logging
import uuid
from pathlib import Path

import boto3
import httpx
from botocore.client import Config

import db
import rag_queue
from chunker import chunk_text
from config import settings

log = logging.getLogger(__name__)

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
_PDF_EXTS = {".pdf"}
_WORD_EXTS = {".docx"}
_PPTX_EXTS = {".pptx"}
_EXCEL_EXTS = {".xlsx", ".xls"}
_TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".xml", ".html"}

_EXT_MIME = {
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".png": "image/png", ".gif": "image/gif",
    ".webp": "image/webp", ".bmp": "image/bmp",
}


def _s3_client():
    return boto3.client(
        "s3",
        endpoint_url=settings.internal_s3_endpoint_url,
        aws_access_key_id=settings.aws_access_key_id,
        aws_secret_access_key=settings.aws_secret_access_key,
        region_name=settings.region_name,
        config=Config(signature_version="s3v4"),
    )


def _download(s3_url: str) -> bytes:
    _, _, rest = s3_url.partition("s3://")
    bucket, _, key = rest.partition("/")
    obj = _s3_client().get_object(Bucket=bucket, Key=key)
    return obj["Body"].read()


def _extract_text(data: bytes, ext: str) -> str:
    if ext in _TEXT_EXTS:
        return data.decode("utf-8", errors="replace")

    if ext in _PDF_EXTS:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    if ext in _WORD_EXTS:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext in _PPTX_EXTS:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)

    if ext in _EXCEL_EXTS:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        rows: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text:
                    rows.append(row_text)
        return "\n".join(rows)

    return ""


async def _describe_image(data: bytes, ext: str) -> str:
    b64 = base64.b64encode(data).decode()
    mime = _EXT_MIME.get(ext, "image/jpeg")
    prompt = "Describe this image in detail for document search indexing."

    if settings.vision_provider == "anthropic":
        import anthropic
        client = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)
        msg = await client.messages.create(
            model=settings.vision_model,
            max_tokens=1024,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64}},
                    {"type": "text", "text": prompt},
                ],
            }],
        )
        return msg.content[0].text
    else:
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(
                f"{settings.ollama_base_url}/api/chat",
                json={
                    "model": settings.vision_model,
                    "messages": [{"role": "user", "content": prompt, "images": [b64]}],
                    "stream": False,
                },
            )
            resp.raise_for_status()
            return resp.json()["message"]["content"]


async def _embed(texts: list[str]) -> list[list[float]]:
    async with httpx.AsyncClient(timeout=120) as client:
        results = await asyncio.gather(*[
            client.post(
                f"{settings.ollama_base_url}/api/embeddings",
                json={"model": settings.embedding_model, "prompt": text},
            )
            for text in texts
        ])
    for r in results:
        r.raise_for_status()
    return [r.json()["embedding"] for r in results]


async def _get_s3_url(doc_id: uuid.UUID) -> str:
    pool = await db.get_pool()
    async with pool.acquire() as conn:
        row = await conn.fetchrow(
            "SELECT s3_url FROM rag.documents WHERE id=$1", doc_id
        )
    if not row:
        raise ValueError(f"Document {doc_id} not found")
    return row["s3_url"]


async def process_document(doc_id: uuid.UUID, s3_url: str) -> None:
    await db.update_document_status(doc_id, "processing")
    try:
        data = _download(s3_url)
        ext = Path(s3_url.rstrip("/").split("/")[-1]).suffix.lower()

        if ext in _IMAGE_EXTS:
            text = await _describe_image(data, ext)
        else:
            text = _extract_text(data, ext)

        if not text.strip():
            raise ValueError(f"No text extracted (ext={ext})")

        chunks = chunk_text(text)
        if not chunks:
            raise ValueError("Chunking produced no output")

        embeddings = await _embed(chunks)
        await db.insert_chunks(doc_id, chunks, embeddings)
        await db.update_document_status(doc_id, "completed", chunk_count=len(chunks))
        log.info("Processed doc %s: %d chunks", doc_id, len(chunks))
    except Exception as exc:
        log.error("Failed to process doc %s: %s", doc_id, exc)
        await db.update_document_status(doc_id, "failed", error=str(exc))


async def worker_loop() -> None:
    log.info("RAG worker started")
    while True:
        try:
            job_id = await rag_queue.dequeue_timeout(timeout=5)
            if job_id is None:
                continue
            doc_id = uuid.UUID(job_id)
            s3_url = await _get_s3_url(doc_id)
            await process_document(doc_id, s3_url)
        except asyncio.CancelledError:
            log.info("RAG worker stopped")
            break
        except Exception as exc:
            log.error("Worker loop error: %s", exc)
            await asyncio.sleep(1)
